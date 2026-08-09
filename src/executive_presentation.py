"""Mayasari Bakery — Executive Management Presentation."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

from src.contracts.paths import (
    CUSTOMER_OPPORTUNITY_DATA,
    CUSTOMER_PERFORMANCE_DATA,
    EXECUTIVE_KPIS_DATA,
    MONTHLY_PERFORMANCE_DATA,
    PRODUCT_PERFORMANCE_DATA,
    PROFITABILITY_SUMMARY_DATA,
)

PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT = PROJECT_ROOT / "reports" / "presentation" / "mayasari_bakery_management_deck.pptx"

FIGURES = PROJECT_ROOT / "reports" / "figures"


def format_currency(value: float) -> str:
    if abs(value) >= 1_000_000_000:
        return f"Rp {value / 1_000_000_000:.1f}B"
    if abs(value) >= 1_000_000:
        return f"Rp {value / 1_000_000:.1f}M"
    if abs(value) >= 1_000:
        return f"Rp {value / 1_000:.1f}K"
    return f"Rp {value:,.0f}"


def load_data() -> dict[str, pd.DataFrame]:
    paths = {
        "executive": EXECUTIVE_KPIS_DATA,
        "profitability": PROFITABILITY_SUMMARY_DATA,
        "monthly": MONTHLY_PERFORMANCE_DATA,
        "customer": CUSTOMER_PERFORMANCE_DATA,
        "product": PRODUCT_PERFORMANCE_DATA,
        "opportunity": CUSTOMER_OPPORTUNITY_DATA,
    }

    for name, path in paths.items():
        if not path.exists():
            raise FileNotFoundError(
                f"{name.title()} dataset not found: {path}"
            )

    return {
        name: pd.read_parquet(path)
        for name, path in paths.items()
    }


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(0.35),
        Inches(12.1),
        Inches(0.7),
    )

    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.62),
            Inches(1.0),
            Inches(11.8),
            Inches(0.45),
        )
        paragraph = subtitle_box.text_frame.paragraphs[0]
        paragraph.text = subtitle
        paragraph.font.size = Pt(12)


def add_bullets(slide, items: list[str]) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(1.7),
        Inches(11.5),
        Inches(5.2),
    )

    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    for index, item in enumerate(items):
        paragraph = (
            frame.paragraphs[0]
            if index == 0
            else frame.add_paragraph()
        )
        paragraph.text = item
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(12)


def add_kpi_cards(
    slide,
    cards: list[tuple[str, str]],
) -> None:
    left = 0.6
    top = 1.7
    width = 2.9
    height = 1.35
    gap = 0.2

    for index, (label, value) in enumerate(cards):
        x = left + index * (width + gap)

        shape = slide.shapes.add_textbox(
            Inches(x),
            Inches(top),
            Inches(width),
            Inches(height),
        )

        frame = shape.text_frame
        frame.clear()

        value_paragraph = frame.paragraphs[0]
        value_paragraph.text = value
        value_paragraph.font.size = Pt(23)
        value_paragraph.font.bold = True

        label_paragraph = frame.add_paragraph()
        label_paragraph.text = label
        label_paragraph.font.size = Pt(11)


def add_image_slide(
    presentation: Presentation,
    title: str,
    image_path: Path,
    subtitle: str | None = None,
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_title(slide, title, subtitle)

    if not image_path.exists():
        raise FileNotFoundError(
            f"Presentation image not found: {image_path}"
        )

    slide.shapes.add_picture(
        str(image_path),
        Inches(0.7),
        Inches(1.45),
        width=Inches(11.9),
    )


def build_presentation(data: dict[str, pd.DataFrame]) -> Presentation:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    executive = data["executive"].iloc[0]
    profitability = data["profitability"].iloc[0]
    monthly = data["monthly"].sort_values("sales_month")
    customer = data["customer"]
    product = data["product"]
    opportunity = data["opportunity"]

    revenue = float(executive["revenue"])
    gross_profit = float(executive["gross_profit"])
    operating_profit = float(executive["operating_profit"])
    operating_expense = float(executive["operating_expense"])

    gross_margin = float(
        profitability["gross_margin_pct"]
    )

    operating_margin = (
        operating_profit / revenue * 100
    )

    # --------------------------------------------------
    # 1. Title
    # --------------------------------------------------

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_title(
        slide,
        "Mayasari Bakery",
        "Executive Management Review — 2025",
    )

    add_bullets(
        slide,
        [
            "Business performance and profitability overview",
            "Customer value and opportunity analysis",
            "Product mix and revenue performance",
            "Management priorities for the next cycle",
        ],
    )

    # --------------------------------------------------
    # 2. Executive Summary
    # --------------------------------------------------

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_title(
        slide,
        "Executive Summary",
        "Current business position",
    )

    add_kpi_cards(
        slide,
        [
            ("Revenue", format_currency(revenue)),
            ("Gross Profit", format_currency(gross_profit)),
            ("Gross Margin", f"{gross_margin:.2f}%"),
            ("Operating Profit", format_currency(operating_profit)),
        ],
    )

    add_bullets(
        slide,
        [
            "The business generated positive operating profit.",
            "Gross margin remains above 40%, providing a meaningful contribution base.",
            "Revenue performance shows meaningful month-to-month variation.",
            "Customer and product concentration should be managed alongside growth.",
        ],
    )

    # --------------------------------------------------
    # 3. Financial Performance
    # --------------------------------------------------

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_title(
        slide,
        "Financial Performance",
        "Profitability structure",
    )

    add_kpi_cards(
        slide,
        [
            ("Revenue", format_currency(revenue)),
            ("Gross Profit", format_currency(gross_profit)),
            ("Operating Expense", format_currency(operating_expense)),
            ("Operating Margin", f"{operating_margin:.2f}%"),
        ],
    )

    add_bullets(
        slide,
        [
            f"Gross profit contribution: {format_currency(gross_profit)}.",
            f"Operating expenses: {format_currency(operating_expense)}.",
            f"Operating profit: {format_currency(operating_profit)}.",
            "Priority: protect margin while pursuing sustainable revenue growth.",
        ],
    )

    # --------------------------------------------------
    # 4. Revenue
    # --------------------------------------------------

    add_image_slide(
        presentation,
        "Revenue Performance",
        FIGURES / "sales" / "revenue_dashboard.png",
        "Monthly revenue trend and performance",
    )

    peak = monthly.loc[monthly["revenue"].idxmax()]
    lowest = monthly.loc[monthly["revenue"].idxmin()]

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_title(
        slide,
        "Revenue Management Takeaways",
        "What management should investigate",
    )

    add_bullets(
        slide,
        [
            f"Peak revenue: {peak['sales_month']} at {format_currency(float(peak['revenue']))}.",
            f"Lowest revenue: {lowest['sales_month']} at {format_currency(float(lowest['revenue']))}.",
            "December showed the strongest monthly growth at approximately +26.45%.",
            "October recorded the weakest month-over-month movement at approximately -7.86%.",
            "Investigate campaign, product availability, and customer activation drivers behind major movements.",
        ],
    )

    # --------------------------------------------------
    # 5. Product
    # --------------------------------------------------

    add_image_slide(
        presentation,
        "Product Performance",
        FIGURES / "products" / "product_dashboard.png",
        "Revenue contribution and product mix",
    )

    top_product = (
        product.sort_values("revenue", ascending=False)
        .iloc[0]
    )

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_title(
        slide,
        "Product Management Takeaways",
        "Revenue versus margin",
    )

    add_bullets(
        slide,
        [
            f"Leading product: {top_product['product_name']} with {format_currency(float(top_product['revenue']))} revenue.",
            "High-revenue products should remain operational priorities.",
            "Products with stronger margins should be evaluated for bundling and cross-selling.",
            "Lower-margin products should be reviewed for pricing, cost, or product-mix opportunities.",
        ],
    )

    # --------------------------------------------------
    # 6. Customer
    # --------------------------------------------------

    add_image_slide(
        presentation,
        "Customer Performance",
        FIGURES / "customers" / "customer_dashboard.png",
        "Customer economics and value distribution",
    )

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_title(
        slide,
        "Customer Value Takeaways",
        "CLV-based management priorities",
    )

    total_clv = customer["annualized_clv"].sum()

    customer_sorted = customer.sort_values(
        "annualized_clv"
    ).reset_index(drop=True)

    customer_sorted["clv_tier"] = pd.qcut(
        customer_sorted["annualized_clv"].rank(
            method="first"
        ),
        q=4,
        labels=[
            "Bronze",
            "Silver",
            "Gold",
            "Platinum",
        ],
    )

    platinum_clv = customer_sorted.loc[
        customer_sorted["clv_tier"] == "Platinum",
        "annualized_clv",
    ].sum()

    platinum_share = platinum_clv / total_clv * 100

    add_bullets(
        slide,
        [
            f"Customer base: {len(customer):,} customers.",
            f"Platinum customers contribute approximately {platinum_share:.2f}% of annualized CLV.",
            "Retention should prioritize economically important customers.",
            "Customer value should be managed using CLV together with observed behavior.",
        ],
    )

    # --------------------------------------------------
    # 7. Opportunity
    # --------------------------------------------------

    add_image_slide(
        presentation,
        "Customer Opportunity",
        FIGURES / "customers" / "customer_dashboard.png",
        "Opportunity prioritization",
    )

    opportunity_summary = (
        opportunity.groupby("opportunity")
        .agg(
            customers=("customer_id", "count"),
            annualized_clv=("annualized_clv", "sum"),
        )
        .sort_values(
            "annualized_clv",
            ascending=False,
        )
    )

    top_opportunities = opportunity_summary.head(3)

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_title(
        slide,
        "Opportunity Priorities",
        "Where management attention should be concentrated",
    )

    bullets = []

    for name, row in top_opportunities.iterrows():
        share = (
            float(row["annualized_clv"])
            / float(opportunity["annualized_clv"].sum())
            * 100
        )

        bullets.append(
            f"{name}: {int(row['customers'])} customers, "
            f"{share:.2f}% of annualized CLV."
        )

    bullets.extend(
        [
            "Rescue customers warrant economically targeted recovery actions.",
            "Review customers require deeper diagnosis before intervention.",
            "Develop customers provide a controlled growth opportunity.",
        ]
    )

    add_bullets(slide, bullets)

    # --------------------------------------------------
    # 8. Profitability
    # --------------------------------------------------

    add_image_slide(
        presentation,
        "Profitability",
        FIGURES / "profitability" / "profitability_dashboard.png",
        "Gross margin and operating profitability",
    )

    # --------------------------------------------------
    # 9. Risks
    # --------------------------------------------------

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_title(
        slide,
        "Key Management Risks",
        "Areas requiring continued monitoring",
    )

    add_bullets(
        slide,
        [
            "Revenue volatility across months can affect planning consistency.",
            "Revenue concentration among leading products creates product-mix dependency.",
            "High-value customer relationships require active retention management.",
            "Margin differences across products can dilute profitability if mix is unmanaged.",
        ],
    )

    # --------------------------------------------------
    # 10. Recommendations
    # --------------------------------------------------

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_title(
        slide,
        "Management Recommendations",
        "Priority actions",
    )

    add_bullets(
        slide,
        [
            "1. Protect revenue momentum — identify repeatable drivers behind strong months.",
            "2. Optimize product mix — balance revenue contribution with gross margin.",
            "3. Protect high-value customers — prioritize Platinum and economically important relationships.",
            "4. Prioritize Rescue and Review — allocate intervention based on customer economics.",
            "5. Protect profitability — ensure revenue growth translates into sustainable operating profit.",
        ],
    )

    # --------------------------------------------------
    # 11. Closing
    # --------------------------------------------------

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_title(
        slide,
        "Management Focus",
        "Recommended decision framework",
    )

    add_bullets(
        slide,
        [
            "Grow revenue with discipline.",
            "Protect gross margin.",
            "Retain economically valuable customers.",
            "Improve product-mix contribution.",
            "Measure campaign and intervention outcomes.",
        ],
    )

    return presentation


def main() -> None:
    data = load_data()

    presentation = build_presentation(data)

    OUTPUT.parent.mkdir(
        parents=True,
        exist_ok=True,
    )

    presentation.save(OUTPUT)

    print(f"Created: {OUTPUT}")
    print(f"Slides : {len(presentation.slides)}")
    print(f"Bytes  : {OUTPUT.stat().st_size}")


if __name__ == "__main__":
    main()
